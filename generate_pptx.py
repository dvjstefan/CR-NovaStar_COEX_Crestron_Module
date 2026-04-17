from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# Brand colors
DARK_BG    = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
ACCENT     = RGBColor(0x00, 0xA8, 0xFF)   # electric blue
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xDD, 0xEE)
MID_GRAY   = RGBColor(0x88, 0xAA, 0xCC)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank


def bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def txb(slide, text, l, t, w, h, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return box


def accent_bar(slide, y=0.72, width=13.33):
    add_rect(slide, 0, y, width, 0.06, ACCENT)


# ── Slide 1 – Title ─────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
bg(s1)
add_rect(s1, 0, 0, 13.33, 7.5, DARK_BG)
add_rect(s1, 0, 2.8, 13.33, 2.2, RGBColor(0x06, 0x30, 0x55))
accent_bar(s1, 2.75)
accent_bar(s1, 4.95)

txb(s1, "SAOA Consulting", 0.7, 0.35, 12, 0.6, 13, color=ACCENT, bold=True)
txb(s1, "NovaStar MX (COEX)", 0.7, 1.05, 12, 1.0, 42, bold=True)
txb(s1, "SIMPL+ Control Module for Crestron 4-Series", 0.7, 2.0, 12, 0.65, 22, color=LIGHT_GRAY)

txb(s1, "Seamless TCP/IP control of NovaStar LED processors\n"
        "— right inside your Crestron ecosystem.", 0.7, 3.05, 10, 1.1, 22, color=WHITE)

txb(s1, "Release 2.20.1  •  saoa.se  •  info@saoa.se",
    0.7, 6.85, 12, 0.5, 11, color=MID_GRAY)


# ── Slide 2 – The Challenge ──────────────────────────────────────────────────
s2 = prs.slides.add_slide(BLANK)
bg(s2)
accent_bar(s2, 0.0)
add_rect(s2, 0, 0, 4.2, 7.5, RGBColor(0x06, 0x30, 0x55))

txb(s2, "THE\nCHALLENGE", 0.35, 1.4, 3.5, 3.0, 36, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)

points = [
    ("LED Walls are Complex", "NovaStar processors require a proprietary TCP/IP API on port 8001 — not standard AV control."),
    ("Integration Gaps", "Most AV control systems lack native NovaStar support, forcing custom workarounds."),
    ("Operator Friction", "Without proper control, technicians must use separate software to manage presets, brightness and display mode."),
]
y = 1.0
for title, body in points:
    add_rect(s2, 4.5, y, 0.07, 0.9, ACCENT)
    txb(s2, title, 4.75, y,       8.2, 0.38, 15, bold=True)
    txb(s2, body,  4.75, y+0.38,  8.2, 0.55, 13, color=LIGHT_GRAY)
    y += 1.55


# ── Slide 3 – The Solution ───────────────────────────────────────────────────
s3 = prs.slides.add_slide(BLANK)
bg(s3)
accent_bar(s3)
add_rect(s3, 0, 0, 13.33, 0.72, ACCENT)
txb(s3, "THE SOLUTION", 0.5, 0.12, 12, 0.55, 22, bold=True, color=DARK_BG)

txb(s3, "NovaStar MX SIMPL+ Module", 0.5, 0.95, 12, 0.75, 32, bold=True)
txb(s3, "A ready-to-use SIMPL+ module that bridges Crestron 4-Series processors\n"
        "directly to NovaStar MX30 LED controllers over TCP/IP.",
    0.5, 1.65, 12, 0.8, 17, color=LIGHT_GRAY)

cols = [
    ("One-click\nDeploy", "Drop the module into any\nSIMPL Windows project."),
    ("Full API\nCoverage", "Presets, brightness,\ndisplay mode & feedback."),
    ("Real-time\nFeedback", "Live status back into\nyour Crestron UI."),
    ("Zero Custom\nCode", "No C# or raw TCP scripting\nrequired."),
]
x = 0.5
for head, body in cols:
    add_rect(s3, x, 2.7, 2.85, 3.8, RGBColor(0x06, 0x30, 0x55))
    add_rect(s3, x, 2.7, 2.85, 0.07, ACCENT)
    txb(s3, head, x+0.18, 2.85, 2.5, 0.9, 18, bold=True, color=ACCENT)
    txb(s3, body, x+0.18, 3.75, 2.5, 1.5, 13, color=LIGHT_GRAY)
    x += 3.1


# ── Slide 4 – Key Features ───────────────────────────────────────────────────
s4 = prs.slides.add_slide(BLANK)
bg(s4)
add_rect(s4, 0, 0, 13.33, 0.72, ACCENT)
txb(s4, "KEY FEATURES", 0.5, 0.12, 12, 0.55, 22, bold=True, color=DARK_BG)

features = [
    ("Preset Recall",        "Instantly recall up to 10 named presets stored on the NovaStar controller via a single analog signal."),
    ("Brightness Control",   "Smooth 0–100 analog brightness control with live feedback to the touch panel."),
    ("Picture Mute/UnMute",  "One-touch mute of the LED wall — ideal for presentations and AV events."),
    ("Display Mode Query",   "Request and display the current operating mode as a readable string feedback."),
    ("TCP Auto-Reconnect",   "Built-in heartbeat and recovery ensures the connection stays alive unattended."),
    ("Preset Name Strings",  "Fetches and exposes all preset names (1–10) as string signals for dynamic UI labels."),
]

col1 = features[:3]
col2 = features[3:]

def feat_block(slide, items, x_start):
    y = 1.0
    for title, desc in items:
        add_rect(slide, x_start, y+0.05, 0.06, 0.65, ACCENT)
        txb(slide, title, x_start+0.22, y,      5.8, 0.38, 15, bold=True)
        txb(slide, desc,  x_start+0.22, y+0.38, 5.8, 0.55, 12, color=LIGHT_GRAY)
        y += 1.7

feat_block(s4, col1, 0.5)
feat_block(s4, col2, 7.0)


# ── Slide 5 – Signal Reference ───────────────────────────────────────────────
s5 = prs.slides.add_slide(BLANK)
bg(s5)
add_rect(s5, 0, 0, 13.33, 0.72, ACCENT)
txb(s5, "SIGNAL REFERENCE", 0.5, 0.12, 12, 0.55, 22, bold=True, color=DARK_BG)

txb(s5, "Inputs", 0.5, 0.85, 6, 0.45, 16, bold=True, color=ACCENT)
txb(s5, "Outputs / Feedback", 7.0, 0.85, 6, 0.45, 16, bold=True, color=ACCENT)

inputs = [
    ("Controller_IP$",    "String",  "IP address of NovaStar controller"),
    ("Controller_Port",   "Analog",  "TCP port (default 8001)"),
    ("Get_Display_Mode",  "Digital", "Request current display mode"),
    ("Picture_Mute",      "Digital", "Mute LED wall output"),
    ("Picture_UnMute",    "Digital", "Unmute LED wall output"),
    ("Get_Presets",       "Digital", "Fetch preset names"),
    ("Set_Brightness",    "Analog",  "Brightness 0–100"),
    ("Set_Preset",        "Analog",  "Recall preset 1–10"),
]
outputs = [
    ("DisplayMode_FB$",   "String",  "Current display mode text"),
    ("Brightness_FB",     "Analog",  "Live brightness value"),
    ("PresetActive",      "Analog",  "Currently active preset number"),
    ("PresetName$[1-10]", "String",  "Preset names from controller"),
]

def table_rows(slide, rows, x, y_start, col_widths):
    y = y_start
    for i, (sig, typ, desc) in enumerate(rows):
        row_color = RGBColor(0x06, 0x30, 0x55) if i % 2 == 0 else RGBColor(0x08, 0x22, 0x40)
        add_rect(slide, x, y, sum(col_widths), 0.38, row_color)
        txb(slide, sig,  x+0.1,                             y+0.05, col_widths[0]-0.1, 0.3, 11, bold=True, color=ACCENT)
        txb(slide, typ,  x+col_widths[0],                   y+0.05, col_widths[1],     0.3, 11, color=LIGHT_GRAY, italic=True)
        txb(slide, desc, x+col_widths[0]+col_widths[1]+0.1, y+0.05, col_widths[2]-0.1, 0.3, 11, color=WHITE)
        y += 0.39
    return y

# Headers
for x_off, label in [(0.5, "Signal"), (2.5, "Type"), (3.7, "Description"),
                      (7.0, "Signal"), (9.0, "Type"), (10.2, "Description")]:
    txb(s5, label, x_off, 1.3, 2.5, 0.3, 11, bold=True, color=MID_GRAY)
add_rect(s5, 0.5, 1.6,  6.0, 0.04, ACCENT)
add_rect(s5, 7.0, 1.6,  6.0, 0.04, ACCENT)

table_rows(s5, inputs,  0.5, 1.65, [2.0, 1.1, 2.9])
table_rows(s5, outputs, 7.0, 1.65, [2.0, 1.1, 2.9])


# ── Slide 6 – System Requirements ───────────────────────────────────────────
s6 = prs.slides.add_slide(BLANK)
bg(s6)
add_rect(s6, 0, 0, 13.33, 0.72, ACCENT)
txb(s6, "SYSTEM REQUIREMENTS", 0.5, 0.12, 12, 0.55, 22, bold=True, color=DARK_BG)

reqs = [
    ("Crestron Processor", "4-Series (CP4, MC4 or equivalent)\nNot designed for 3-Series"),
    ("Firmware",           "Crestron firmware v2.0 or newer"),
    ("SIMPL Windows",      "Version 4.17.00 or later"),
    ("Network",            "TCP port 8001 open between the\nCrestron processor and NovaStar unit"),
    ("NovaStar Hardware",  "MX30 controller (COEX platform)\nFirmware compatible with COEX API"),
]
x = 0.5
for head, body in reqs:
    add_rect(s6, x, 1.1, 2.35, 4.5, RGBColor(0x06, 0x30, 0x55))
    add_rect(s6, x, 1.1, 2.35, 0.07, ACCENT)
    txb(s6, head, x+0.15, 1.25, 2.05, 0.65, 14, bold=True, color=ACCENT)
    txb(s6, body, x+0.15, 1.9,  2.05, 2.5,  12, color=LIGHT_GRAY)
    x += 2.57

txb(s6, "This module is a drop-in SIMPL+ component — no additional middleware, servers or drivers required.",
    0.5, 5.9, 12, 0.5, 13, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ── Slide 7 – Why SAOA ───────────────────────────────────────────────────────
s7 = prs.slides.add_slide(BLANK)
bg(s7)
add_rect(s7, 0, 0, 5.5, 7.5, RGBColor(0x06, 0x30, 0x55))
accent_bar(s7, 0.0, 5.5)

txb(s7, "WHY\nSAOA\nCONSULTING", 0.5, 1.3, 4.5, 3.5, 36, bold=True, color=ACCENT)
txb(s7, "Your AV integration\nspecialist.", 0.5, 4.8, 4.5, 1.5, 18, color=LIGHT_GRAY)

bullets = [
    ("Crestron Certified", "Deep expertise in SIMPL, SIMPL# and 4-Series architecture."),
    ("NovaStar Specialists", "Direct experience with COEX API integration across live events and permanent installs."),
    ("Full Lifecycle Support", "From design and programming to commissioning and on-site support."),
    ("Custom Development", "Need a tailored module or full system design? We deliver."),
]
y = 0.9
for title, body in bullets:
    add_rect(s7, 5.8, y, 0.07, 0.85, ACCENT)
    txb(s7, title, 6.0, y,       7.0, 0.38, 15, bold=True)
    txb(s7, body,  6.0, y+0.38,  7.0, 0.55, 13, color=LIGHT_GRAY)
    y += 1.55


# ── Slide 8 – Track Record ───────────────────────────────────────────────────
s8 = prs.slides.add_slide(BLANK)
bg(s8)
add_rect(s8, 0, 0, 13.33, 0.72, ACCENT)
txb(s8, "PROVEN TRACK RECORD", 0.5, 0.12, 12, 0.55, 22, bold=True, color=DARK_BG)

releases = [
    ("v2.1",   "Nov 2025", "Initial public release of the NovaStar MX COEX module."),
    ("v2.20",  "Jan 2026", "Major stability update: improved MX30 state handling, heartbeat monitoring and connection recovery."),
    ("v2.20.1","Apr 2026", "Packaging cleanup: removed unused dependencies, reduced compiled .clz size from 3.8 MB to 1.4 MB."),
]

y = 1.1
for ver, date, notes in releases:
    add_rect(s8, 0.5,  y,      1.4,  1.0, RGBColor(0x00, 0x68, 0xAA))
    add_rect(s8, 1.9,  y,      10.9, 1.0, RGBColor(0x06, 0x30, 0x55))
    add_rect(s8, 0.5,  y,      12.4, 0.06, ACCENT)
    txb(s8, ver,   0.55, y+0.1, 1.3, 0.45, 20, bold=True, color=WHITE,   align=PP_ALIGN.CENTER)
    txb(s8, date,  0.55, y+0.55,1.3, 0.35, 11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    txb(s8, notes, 2.05, y+0.2, 10.5,0.65, 13, color=WHITE)
    y += 1.25

txb(s8, "Actively maintained with a clear versioning policy and backward-compatible upgrades.",
    0.5, 5.2, 12, 0.5, 14, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ── Slide 9 – Call to Action ─────────────────────────────────────────────────
s9 = prs.slides.add_slide(BLANK)
bg(s9)
add_rect(s9, 0, 0, 13.33, 7.5, DARK_BG)
add_rect(s9, 0, 2.5, 13.33, 2.8, RGBColor(0x06, 0x30, 0x55))
accent_bar(s9, 2.45)
accent_bar(s9, 5.25)

txb(s9, "READY TO INTEGRATE?", 0.7, 0.6, 12, 0.8, 38, bold=True, align=PP_ALIGN.CENTER)
txb(s9, "Get the NovaStar MX SIMPL+ module working in your project today.",
    0.7, 1.45, 12, 0.6, 18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

txb(s9, "Contact Us", 0.7, 2.7, 12, 0.55, 22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
txb(s9, "SAOA Consulting", 0.7, 3.3, 12, 0.55, 20, bold=True, align=PP_ALIGN.CENTER)
txb(s9, "info@saoa.se", 0.7, 3.9, 12, 0.5, 18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
txb(s9, "https://saoa.se", 0.7, 4.4, 12, 0.5, 18, color=ACCENT, align=PP_ALIGN.CENTER)

txb(s9, "Release 2.20.1  •  Crestron 4-Series  •  TCP/IP  •  NovaStar COEX API",
    0.7, 6.85, 12, 0.5, 11, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────────────────
out = "/home/user/CR-NovaStar_COEX_Crestron_Module/NovaStar_MX_SalesPresentation.pptx"
prs.save(out)
print(f"Saved: {out}")
